#rag_service
from pinecone import Pinecone
from services.vector_store import retrieve_from_kb
from services.hf_model import ask_gpt
import re
import asyncio
import os
import tempfile
import requests
from typing import List
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_pinecone import PineconeVectorStore
from urllib.parse import urlparse
from config.settings import settings


# Additional imports for new file formats
from openpyxl import load_workbook
from PIL import Image
import pytesseract
import subprocess
from pptx import Presentation
from pdf2image import convert_from_path
from docx import Document as DocxDocument

# Initialize Pinecone with error handling
try:
    pc = Pinecone(api_key=settings.PINECONE_API_KEY)
    index = pc.Index(settings.PINECONE_INDEX_NAME)
except Exception as e:
    print(f"Error initializing Pinecone: {e}")
    raise

def generate_namespace_from_url(url: str) -> str:
    """Generate namespace matching the embedding script logic"""
    try:
        # Parse URL to get the filename part
        parsed = urlparse(url)
        filename = parsed.path.split('/')[-1]
        
        # Remove query parameters if they're part of the filename
        if '?' in filename:
            filename = filename.split('?')[0]
            
        # Remove extension (matching your embedding script logic)
        name_without_ext = os.path.splitext(filename)[0]
        
        # Replace non-alphanumeric characters with underscores and convert to lowercase
        # This matches exactly: re.sub(r'[^a-zA-Z0-9]+', '_', name_without_ext).strip('_').lower()
        namespace = re.sub(r'[^a-zA-Z0-9]+', '_', name_without_ext).strip('_').lower()
        
        # Ensure namespace is not empty (matching your embedding script)
        if not namespace:
            namespace = "default_namespace"
            
        return namespace
        
    except Exception as e:
        print(f"Error generating namespace: {e}")
        return "default_namespace"

def clean_metadata_for_pinecone(metadata: dict, max_total_size=2000) -> dict:
    """
    Clean metadata to ensure it's compatible with Pinecone and under size limit.
    """
    cleaned = {}
    total_size = 0

    for key, value in metadata.items():
        if value is None:
            value = "unknown"
        elif isinstance(value, (list, dict)):
            value = str(value)
        elif not isinstance(value, (str, int, float, bool)):
            value = str(value)

        value_str = str(value)

        if len(value_str) > 300:  # Truncate very long values
            value_str = value_str[:300] + "..."

        total_size += len(key) + len(value_str)
        if total_size > max_total_size:
            print(f"⚠️ Skipping metadata key '{key}' to stay under limit")
            continue

        cleaned[key] = value_str

    return cleaned

# New file format extraction functions
def extract_text_from_xlsx(file_path: str) -> list[str]:
    """Extract text from Excel files"""
    wb = load_workbook(filename=file_path, read_only=True, data_only=True)
    chunks = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        # Step 1: Identify the header row
        header_row_index = None
        for i, row in enumerate(rows):
            if row and {"Name", "Mobile Number", "Pincode", "Salary"}.issubset(set(str(cell) for cell in row if cell)):
                header_row_index = i
                break
        # Step 2: Extract message-like unstructured content before the table
        for row in rows[:header_row_index]:
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                chunks.append(row_text.strip())
        # Step 3: Extract structured data after the header
        if header_row_index is not None:
            headers = [str(cell) for cell in rows[header_row_index] if cell]
            for row in rows[header_row_index + 1:]:
                if any(row):  # skip empty rows
                    row_text = " | ".join(
                        f"{header}: {cell}" for header, cell in zip(headers, row) if cell is not None
                    )
                    if row_text.strip():
                        chunks.append(row_text.strip())
    print(f"📊 Extracted {len(chunks)} chunks from Excel")
    print(f"🔍 First chunk (100 chars): {repr(chunks[0][:100]) if chunks else 'No chunks'}")
    return chunks

def extract_text_from_image(image_path: str) -> list[str]:
    """Extract text from images using OCR"""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    
    print(f"🖼️ Extracted text from image: {text[:100]}...")  # Log first 100 chars for debugging
    chunks = [chunk.strip() for chunk in text.split("\n\n") if chunk.strip()]
    print(f"🖼️ Extracted {len(chunks)} chunks from image")
    print(f"🔍 First chunk (100 chars): {repr(chunks[0][:100]) if chunks else 'No chunks'}")
    return chunks

def pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """
    Converts a PPTX file to PDF using LibreOffice and returns the actual PDF path.
    """
    command = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_path
    ]
    try:
        result = subprocess.run(command, capture_output=True, check=True)
        print("✅ LibreOffice Output:", result.stdout.decode())
    except subprocess.CalledProcessError as e:
        print("❌ LibreOffice failed:", e.stderr.decode())
        raise RuntimeError("PDF conversion failed") from e
    pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_filename)
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not found at expected location: {pdf_path}")
    return pdf_path

def extract_text_from_pdf_ocr(pdf_path: str) -> list[str]:
    """Extract text from PDF using OCR"""
    slides = convert_from_path(pdf_path, dpi=300)
    all_chunks = []
    for i, image in enumerate(slides):
        print(f"🔍 OCR on slide {i+1}")
        text = pytesseract.image_to_string(image)
        chunks = [chunk.strip() for chunk in text.split("\n\n") if chunk.strip()]
        all_chunks.extend(chunks)
    print(f"🖼️ Extracted {len(all_chunks)} chunks from OCR (PDF-based)")
    return all_chunks

def extract_text_from_pptx(pptx_path: str) -> list[str]:
    """Extract text from PowerPoint files"""
    prs = Presentation(pptx_path)
    full_text = []
    slides_needing_ocr = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            full_text.append("\n".join(slide_text))
        else:
            print(f"⚠️ No text found on slide {i+1}, marking for OCR.")
            slides_needing_ocr.append(i)
    if not slides_needing_ocr:
        print("📊 All text extracted from PPTX without OCR.")
        return [chunk.strip() for chunk in full_text if chunk.strip()]
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = pptx_to_pdf(pptx_path, tmpdir)
        all_images = convert_from_path(pdf_path, dpi=300)
        for i in slides_needing_ocr:
            image = all_images[i]
            print(f"🔍 OCR on slide {i+1}")
            ocr_text = pytesseract.image_to_string(image)
            ocr_chunks = [chunk.strip() for chunk in ocr_text.split("\n\n") if chunk.strip()]
            full_text.extend(ocr_chunks)
    print(f"📊 Extracted {len(full_text)} chunks from PPTX (with OCR fallback)")
    print(f"🔍 First chunk (100 chars): {repr(full_text[0][:100]) if full_text else 'No chunks'}")
    return [chunk.strip() for chunk in full_text if chunk.strip()]

def extract_text_from_txt(txt_path: str) -> list[str]:
    """Extract text from plain text files"""
    with open(txt_path, "r", encoding="utf-8") as f:
        content = f.read()
    chunks = [chunk.strip() for chunk in content.split("\n\n") if chunk.strip()]
    print(f"📜 Extracted {len(chunks)} chunks from TXT")
    print(f"🔍 First chunk (100 chars): {repr(chunks[0][:100]) if chunks else 'No chunks'}")
    return chunks

def extract_text_from_docx(docx_path: str) -> list[str]:
    """Extract text from Word documents"""
    doc = DocxDocument(docx_path)
    full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    
    chunks = [chunk.strip() for chunk in full_text.split("\n\n") if chunk.strip()]
    
    print(f"📝 Extracted {len(chunks)} chunks from DOCX")
    print(f"🔍 First chunk (100 chars): {repr(chunks[0][:100]) if chunks else 'No chunks'}")
    return chunks

def detect_file_type(file_path: str) -> str:
    """Detect file type based on extension"""
    _, ext = os.path.splitext(file_path.lower())
    return ext[1:]  # Remove the dot

def load_and_extract_text_from_file(file_path: str) -> List[str]:
    """
    Load and extract text from various file formats.
    Returns a list of text chunks.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_name = os.path.basename(file_path)
    file_type = detect_file_type(file_path)
    
    print(f"📄 Loading file: {file_name} (type: {file_type})")
    
    try:
        if file_type == 'pdf':
            # Use existing PDF loading logic
            loader = PyMuPDFLoader(file_path)
            docs = loader.load()
            if not docs:
                raise ValueError(f"No content found in PDF: {file_name}")
            
            # Convert to text chunks
            chunks = []
            for doc in docs:
                if doc.page_content.strip():
                    chunks.append(doc.page_content.strip())
            print(f"📖 Loaded {len(chunks)} chunks from PDF")
            
        elif file_type in ['xlsx', 'xls']:
            chunks = extract_text_from_xlsx(file_path)
            
        elif file_type in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif']:
            chunks = extract_text_from_image(file_path)
            
        elif file_type in ['pptx', 'ppt']:
            chunks = extract_text_from_pptx(file_path)
            
        elif file_type == 'txt':
            chunks = extract_text_from_txt(file_path)
            
        elif file_type in ['docx', 'doc']:
            chunks = extract_text_from_docx(file_path)
            
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        if not chunks:
            raise ValueError(f"No text content extracted from file: {file_name}")
        
        return chunks
        
    except Exception as e:
        raise Exception(f"Failed to extract text from {file_name}: {str(e)}")

def load_and_split_pdf(file_path: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[Document]:
    """Load PDF and split into chunks using RecursiveCharacterTextSplitter."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_name = os.path.basename(file_path)
    print(f"📄 Loading PDF: {file_name}")

    # Load PDF using PyMuPDFLoader
    try:
        loader = PyMuPDFLoader(file_path)
        docs = loader.load()
    except Exception as e:
        raise Exception(f"Failed to load PDF {file_name}: {str(e)}")

    if not docs:
        raise ValueError(f"No content found in PDF: {file_name}")

    print(f"📖 Loaded {len(docs)} pages from PDF")

    # Add source file metadata
    for doc in docs:
        doc.metadata["source_file"] = file_name
        if "page" not in doc.metadata:
            doc.metadata["page"] = "unknown"

    # Split documents using RecursiveCharacterTextSplitter
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    chunks = text_splitter.split_documents(docs)
    print(f"📚 Split into {len(chunks)} chunks")
    
    return chunks

def load_and_split_file(file_path: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[Document]:
    """
    Load any supported file format and split into Document chunks.
    This is a generalized version that handles all file types.
    """
    file_name = os.path.basename(file_path)
    file_type = detect_file_type(file_path)
    
    # For PDF files, use the existing detailed PDF processing
    if file_type == 'pdf':
        return load_and_split_pdf(file_path, chunk_size, chunk_overlap)
    
    # For other file types, extract text and convert to Documents
    try:
        text_chunks = load_and_extract_text_from_file(file_path)
        
        # Convert text chunks to Document objects
        documents = []
        for i, chunk in enumerate(text_chunks):
            doc = Document(
                page_content=chunk,
                metadata={
                    "source_file": file_name,
                    "file_type": file_type,
                    "chunk_index": i
                }
            )
            documents.append(doc)
        
        # Apply text splitting if chunks are too large
        if any(len(doc.page_content) > chunk_size for doc in documents):
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            documents = text_splitter.split_documents(documents)
            print(f"📚 Further split into {len(documents)} chunks")
        
        return documents
        
    except Exception as e:
        raise Exception(f"Failed to process {file_name}: {str(e)}")

def download_file_from_url(url: str) -> str:
    """Download file from URL to temporary file and return the file path."""
    try:
        print(f"📥 Downloading file from: {url}")
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        # Try to determine file extension from URL or Content-Type
        parsed_url = urlparse(url)
        url_path = parsed_url.path
        if '.' in url_path:
            file_ext = os.path.splitext(url_path)[1]
        else:
            # Try to determine from Content-Type
            content_type = response.headers.get('content-type', '').lower()
            if 'pdf' in content_type:
                file_ext = '.pdf'
            elif 'excel' in content_type or 'spreadsheet' in content_type:
                file_ext = '.xlsx'
            elif 'powerpoint' in content_type or 'presentation' in content_type:
                file_ext = '.pptx'
            elif 'word' in content_type or 'document' in content_type:
                file_ext = '.docx'
            elif 'text' in content_type:
                file_ext = '.txt'
            elif 'image' in content_type:
                if 'jpeg' in content_type:
                    file_ext = '.jpg'
                elif 'png' in content_type:
                    file_ext = '.png'
                else:
                    file_ext = '.jpg'  # Default for images
            else:
                file_ext = ''  # No extension
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            for chunk in response.iter_content(chunk_size=8192):
                temp_file.write(chunk)
            temp_path = temp_file.name
        
        print(f"✅ File downloaded to: {temp_path}")
        return temp_path
        
    except Exception as e:
        print(f"❌ Error downloading file: {e}")
        raise

def download_pdf_from_url(url: str) -> str:
    """Download PDF from URL to temporary file and return the file path."""
    # This function is kept for backward compatibility
    return download_file_from_url(url)

async def embed_file_to_pinecone(file_url: str, namespace: str, chunk_size: int = 1000, chunk_overlap: int = 200):
    """Process any supported file format from URL and upload chunks to Pinecone."""
    print(f"🚀 Starting file embedding for namespace: {namespace}")
    
    temp_file_path = None
    try:
        # Download file to temporary location
        temp_file_path = download_file_from_url(file_url)
        
        # Load and split file based on its type
        chunks = load_and_split_file(temp_file_path, chunk_size, chunk_overlap)
        
        if not chunks:
            print("❌ No chunks to process")
            return False

        # Initialize embedding model (using local sentence transformer)
        try:
            embedding_model = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )
            print("🤖 Local embedding model initialized")
        except Exception as e:
            print(f"❌ Error initializing embedding model: {str(e)}")
            return False

        # Clean metadata for all chunks
        cleaned_chunks = []
        for chunk in chunks:
            cleaned_metadata = clean_metadata_for_pinecone(chunk.metadata)
            # Add URL to metadata
            cleaned_metadata["source_url"] = file_url
            cleaned_chunk = Document(
                page_content=chunk.page_content, 
                metadata=cleaned_metadata
            )
            cleaned_chunks.append(cleaned_chunk)

        # Upload to Pinecone
        try:
            print("📌 Uploading to Pinecone...")
            
            vectorstore = PineconeVectorStore.from_documents(
                documents=cleaned_chunks,
                embedding=embedding_model,
                index_name=settings.PINECONE_INDEX_NAME,
                namespace=namespace
            )
            
            print(f"✅ Successfully embedded and stored {len(chunks)} chunks")
            print(f"📍 Namespace: {namespace}")
            print(f"🔍 Index: {settings.PINECONE_INDEX_NAME}")
            return True
            
        except Exception as e:
            print(f"❌ Error uploading to Pinecone: {str(e)}")
            return False
            
    except Exception as e:
        print(f"❌ Error in file embedding process: {e}")
        return False
    finally:
        # Clean up temporary file
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
                print(f"🗑️ Cleaned up temporary file: {temp_file_path}")
            except Exception as e:
                print(f"⚠️ Error cleaning up temp file: {e}")

async def embed_pdf_to_pinecone(pdf_url: str, namespace: str, chunk_size: int = 1000, chunk_overlap: int = 200):
    """Process PDF from URL and upload chunks to Pinecone."""
    # This function is kept for backward compatibility, but now uses the generalized embed_file_to_pinecone
    return await embed_file_to_pinecone(pdf_url, namespace, chunk_size, chunk_overlap)

async def list_available_namespaces() -> list[str]:
    """Helper function to list all available namespaces in the Pinecone index"""
    try:
        stats = index.describe_index_stats()
        namespaces = list(stats.namespaces.keys()) if stats.namespaces else []
        return namespaces
    except Exception as e:
        print(f"Error retrieving namespaces: {e}")
        return []

async def process_documents_and_questions(file_url: str, questions: list[str], namespace: str = None) -> dict:
    """
    Process documents (any supported format) and questions.
    Updated to handle multiple file formats while maintaining backward compatibility.
    """
    print(f"Processing questions for file URL: {file_url}")
    
    try:
        # Step 1: Handle namespace determination
        if namespace:
            # Use provided namespace directly
            agent_id = namespace
            print(f"📂 Using provided namespace: '{agent_id}'")
        else:
            # Generate namespace using the same logic as your embedding scripts
            agent_id = generate_namespace_from_url(file_url)
            print(f"📂 Generated namespace: '{agent_id}'")
        
        # Debug: Check what namespaces actually exist
        try:
            stats = index.describe_index_stats()
            existing_namespaces = list(stats.namespaces.keys()) if stats.namespaces else []
            print(f"🔍 Available namespaces: {existing_namespaces}")
            
            namespace_exists = agent_id in existing_namespaces
            
            if not namespace_exists:
                print(f"⚠️ Namespace '{agent_id}' not found in existing namespaces")
                
                if not namespace:  # Only auto-select if namespace wasn't provided
                    # Try common patterns based on your embedding scripts
                    possible_namespaces = [
                        agent_id,  # Direct match
                        "extracted_text_embedding",  # For text files
                        f"{agent_id}_pdf",  # With suffix
                        f"doc_{agent_id}",  # With prefix
                    ]
                    
                    # Also try partial matches for existing namespaces
                    for existing_ns in existing_namespaces:
                        if agent_id in existing_ns.lower() or existing_ns.lower() in agent_id:
                            possible_namespaces.append(existing_ns)
                    
                    found_namespace = None
                    for candidate in possible_namespaces:
                        if candidate in existing_namespaces:
                            found_namespace = candidate
                            break
                    
                    if found_namespace:
                        agent_id = found_namespace
                        print(f"🔄 Found matching namespace: '{agent_id}'")
                        namespace_exists = True
                    else:
                        # Namespace doesn't exist, so embed the file
                        print(f"📥 Namespace not found. Embedding file from URL...")
                        embedding_success = await embed_file_to_pinecone(
                            file_url=file_url,
                            namespace=agent_id,
                            chunk_size=500,
                            chunk_overlap=100
                        )
                        
                        if not embedding_success:
                            raise Exception(f"Failed to embed file from URL: {file_url}")
                        
                        print(f"✅ Successfully created namespace '{agent_id}' with file embeddings")
                        namespace_exists = True
                else:
                    # Provided namespace doesn't exist, embed the file
                    print(f"📥 Provided namespace '{namespace}' not found. Embedding file from URL...")
                    embedding_success = await embed_file_to_pinecone(
                        file_url=file_url,
                        namespace=agent_id,
                        chunk_size=500,
                        chunk_overlap=100
                    )
                    
                    if not embedding_success:
                        raise Exception(f"Failed to embed file for namespace '{namespace}'")
                    
                    print(f"✅ Successfully created namespace '{agent_id}' with file embeddings")
                
        except Exception as e:
            if "Failed to embed file" in str(e):
                raise e
            print(f"Error checking namespaces: {e}")
            # Continue with the generated namespace anyway

        # Step 2: Parallel question processing with reduced concurrency
        semaphore = asyncio.Semaphore(3)  # Reduced from 10 to 3 to avoid rate limits

        async def process_question(index: int, question: str) -> tuple[int, str, str]:
            async with semaphore:
                for attempt in range(3):
                    try:
                        retrieval_input = {"query": question, "agent_id": agent_id, "top_k": 3}
                        retrieved = await retrieve_from_kb(retrieval_input)
                        retrieved_chunks = retrieved.get("chunks", [])
                        
                        if not retrieved_chunks:
                            print(f"⚠️ Q{index}: No chunks retrieved for question: {question[:50]}...")
                            return (index, question, "I couldn't find relevant information to answer this question.")

                        max_context_chars = 3000
                        context = "\n".join(retrieved_chunks)[:max_context_chars]

                        print(f"✏️ Q{index}: Processing question: {question[:50]}...")
                        answer = await ask_gpt(context, question)
                        return (index, question, answer)
                        
                    except Exception as e:
                        print(f"⚠️ Q{index}: Attempt {attempt + 1} failed with error: {e}")
                        if attempt < 2:  # Don't sleep on last attempt
                            await asyncio.sleep(2 ** attempt)  # Exponential backoff
                
                return (index, question, "Sorry, I couldn't find relevant information to answer this question.")

        print(f"🧠 Parallel processing {len(questions)} questions...")
        
        if not questions:
            return {}
            
        # Add timeout for question processing
        try:
            tasks = [asyncio.create_task(process_question(i, q)) for i, q in enumerate(questions)]
            responses = await asyncio.wait_for(asyncio.gather(*tasks), timeout=120)  # 2 minute timeout
        except asyncio.TimeoutError:
            print("⚠️ Question processing timed out")
            raise Exception("Processing timed out. Please try with fewer questions or a smaller document.")

        # Step 3: Return sorted results
        results = {q: ans for _, q, ans in sorted(responses)}
        return results
        
    except Exception as e:
        print(f"❌ Error in process_documents_and_questions: {e}")
        raise